import os
from abc import ABC, abstractmethod
from typing import List


# CLASE BASE INTERFAZ

class LectorArchivo(ABC):
    """
    Clase base abstracta para todos los lectores de archivos.
    
    Define el contrato que deben cumplir todos los lectores:
    - Método leer(): lee el contenido del archivo
    - Método puede_leer(): verifica si puede leer una extensión
    """
    
    @abstractmethod
    def leer(self, ruta: str) -> str:
        """
        Lee el contenido de un archivo y lo retorna como texto en minúsculas.
        
        Args:
            ruta: Ruta completa del archivo
            
        Returns:
            Contenido del archivo en minúsculas (para búsqueda case-insensitive)
        """
        pass
    
    @abstractmethod
    def puede_leer(self, extension: str) -> bool:
        """
        Verifica si este lector puede procesar la extensión dada.
        
        Args:
            extension: Extensión del archivo (ej: '.pdf', '.docx')
            
        Returns:
            True si puede leer, False en caso contrario
        """
        pass


# LECTORES CONCRETOS (uno por cada formato)

class LectorTexto(LectorArchivo):
    """
    Lector para archivos de texto plano.
    Soporta múltiples encodings para manejar diferentes idiomas.
    """
    
    EXTENSIONES = {
        '.txt', '.log', '.csv', '.json', '.xml', '.html', 
        '.py', '.js', '.css', '.md', '.ini', '.conf'
    }
    
    ENCODINGS = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
    
    def puede_leer(self, extension: str) -> bool:
        return extension.lower() in self.EXTENSIONES
    
    def leer(self, ruta: str) -> str:
        """Lee archivo de texto probando múltiples encodings."""
        for encoding in self.ENCODINGS:
            try:
                with open(ruta, 'r', encoding=encoding) as f:
                    return f.read().lower()
            except (UnicodeDecodeError, IOError):
                continue
        return ""


class LectorPDF(LectorArchivo):
    """Lector para archivos PDF usando PyPDF2."""
    
    def puede_leer(self, extension: str) -> bool:
        return extension.lower() == '.pdf'
    
    def leer(self, ruta: str) -> str:
        """Extrae texto de todas las páginas del PDF."""
        try:
            import PyPDF2
            with open(ruta, 'rb') as f:
                pdf = PyPDF2.PdfReader(f)
                texto = ''
                for page in pdf.pages:
                    texto += page.extract_text()
                return texto.lower()
        except Exception as e:
            print(f"Error leyendo PDF {ruta}: {e}")
            return ""


class LectorDOCX(LectorArchivo):
    """Lector para archivos Word (.docx)."""
    
    def puede_leer(self, extension: str) -> bool:
        return extension.lower() == '.docx'
    
    def leer(self, ruta: str) -> str:
        """Extrae texto de todos los párrafos del documento."""
        try:
            import docx
            doc = docx.Document(ruta)
            texto = '\n'.join([p.text for p in doc.paragraphs])
            return texto.lower()
        except Exception as e:
            print(f"Error leyendo DOCX {ruta}: {e}")
            return ""


class LectorXLSX(LectorArchivo):
    """Lector para archivos Excel (.xlsx, .xls)."""
    
    def puede_leer(self, extension: str) -> bool:
        return extension.lower() in {'.xlsx', '.xls'}
    
    def leer(self, ruta: str) -> str:
        """Extrae contenido de todas las hojas y celdas."""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(ruta, read_only=True, data_only=True)
            texto = ''
            
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            texto += str(cell) + ' '
                    texto += '\n'
            
            wb.close()
            return texto.lower()
        except Exception as e:
            print(f"Error leyendo XLSX {ruta}: {e}")
            return ""


class LectorPPTX(LectorArchivo):
    """Lector para presentaciones PowerPoint (.pptx)."""
    
    def puede_leer(self, extension: str) -> bool:
        return extension.lower() == '.pptx'
    
    def leer(self, ruta: str) -> str:
        """Extrae texto de todas las diapositivas."""
        try:
            from pptx import Presentation
            prs = Presentation(ruta)
            texto = ''
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texto += shape.text + '\n'
            
            return texto.lower()
        except Exception as e:
            print(f"Error leyendo PPTX {ruta}: {e}")
            return ""


class LectorGenerico(LectorArchivo):
    """
    Lector genérico para archivos no soportados.
    Actúa como fallback cuando ningún otro lector puede manejar el archivo.
    """
    
    def puede_leer(self, extension: str) -> bool:
        return True  # Acepta cualquier extensión
    
    def leer(self, ruta: str) -> str:
        """Retorna string vacío para archivos no soportados."""
        return ""


# FACTORY (El cerebro del sistema)

class LectorFactory:
    """
    Factory que crea el lector apropiado según la extensión del archivo.
    
        lector = LectorFactory.crear_lector('documento.pdf')
        contenido = lector.leer('documento.pdf')
    """
    
    # Registro centralizado de lectores
    # El orden importa: se prueban en secuencia
    _lectores: List[LectorArchivo] = [
        LectorTexto(),
        LectorPDF(),
        LectorDOCX(),
        LectorXLSX(),
        LectorPPTX(),
        LectorGenerico()  # SIEMPRE al final como fallback
    ]
    
    @classmethod
    def crear_lector(cls, ruta: str) -> LectorArchivo:
        """
        Crea y retorna el lector apropiado para el archivo.
        
        Args:
            ruta: Ruta del archivo a leer
            
        Returns:
            Instancia del lector apropiado
            
        Example:
            >>> lector = LectorFactory.crear_lector('informe.pdf')
            >>> type(lector).__name__
            'LectorPDF'
        """
        extension = os.path.splitext(ruta)[1].lower()
        
        # Buscar el primer lector que pueda manejar esta extensión
        for lector in cls._lectores:
            if lector.puede_leer(extension):
                return lector
        
        # Fallback (nunca debería llegar aquí por LectorGenerico)
        return LectorGenerico()
    
    @classmethod
    def registrar_lector(cls, lector: LectorArchivo, posicion: int = -1):
        """
        Registra un nuevo lector dinámicamente sin modificar el código.
        
        Esto permite extensibilidad total:
        - Puedes agregar soporte para nuevos formatos en runtime
        - No necesitas modificar LectorFactory
        - Perfecto para plugins o extensiones
        
        Args:
            lector: Instancia del nuevo lector
            posicion: Posición en la lista de lectores
                     -1 = antes del LectorGenerico (recomendado)
                     0 = al principio (máxima prioridad)
        
        Example:
            >>> lector_epub = LectorEPUB()
            >>> LectorFactory.registrar_lector(lector_epub)
        """
        if posicion == -1:
            # Insertar antes del LectorGenerico (última posición)
            cls._lectores.insert(len(cls._lectores) - 1, lector)
        else:
            cls._lectores.insert(posicion, lector)
    
    @classmethod
    def obtener_extensiones_soportadas(cls) -> set:
        """
        Retorna todas las extensiones soportadas por los lectores registrados.
        
        Returns:
            Set con todas las extensiones soportadas
            
        Example:
            >>> extensiones = LectorFactory.obtener_extensiones_soportadas()
            >>> '.pdf' in extensiones
            True
        """
        extensiones = set()
        
        for lector in cls._lectores:
            if hasattr(lector, 'EXTENSIONES'):
                extensiones.update(lector.EXTENSIONES)
            elif lector.puede_leer('.test'):  # Test genérico
                continue  # Ignorar lectores genéricos
            else:
                # Para lectores individuales, probar extensiones comunes
                for ext in ['.pdf', '.docx', '.xlsx', '.pptx', '.txt']:
                    if lector.puede_leer(ext):
                        extensiones.add(ext)
        
        return extensiones


# FUNCIÓN DE CONVENIENCIA (para usar en tu código existente)

def leer_contenido_archivo(ruta: str) -> str:
    """
    Lee el contenido de un archivo usando Factory Pattern.
    
    Esta función reemplaza tu función anterior y es mucho más simple.
    
    Args:
        ruta: Ruta completa del archivo a leer
        
    Returns:
        Contenido del archivo en minúsculas, o string vacío si hay error
        
    Example:
        >>> contenido = leer_contenido_archivo('documento.pdf')
        >>> 'inteligencia artificial' in contenido
        True
    """
    try:
        # El Factory decide qué lector usar
        lector = LectorFactory.crear_lector(ruta)
        
        # El lector hace su trabajo
        return lector.leer(ruta)
        
    except Exception as e:
        print(f"Error leyendo archivo {ruta}: {e}")
        return ""


# EJEMPLO DE USO Y PRUEBAS

if __name__ == '__main__':
    """
    Ejemplos de uso del Factory Pattern.
    Ejecuta este archivo directamente para ver cómo funciona.
    """
    
    print("=" * 70)
    print("DEMOSTRACIÓN DEL FACTORY PATTERN")
    print("=" * 70)
     
    # Ejemplo 1: Crear lectores para diferentes archivos
    print("\n1. CREACIÓN DE LECTORES:")
    print("-" * 50)
    
    archivos_prueba = [
        'documento.pdf',
        'informe.docx',
        'datos.xlsx',
        'presentacion.pptx',
        'notas.txt',
        'archivo.desconocido'
    ]
    
    for archivo in archivos_prueba:
        lector = LectorFactory.crear_lector(archivo)
        print(f"  {archivo:25} → {type(lector).__name__}")
    
    # Ejemplo 2: Extensiones soportadas
    print("\n2. EXTENSIONES SOPORTADAS:")
    print("-" * 50)
    extensiones = LectorFactory.obtener_extensiones_soportadas()
    print(f"  Total: {len(extensiones)} extensiones")
    print(f"  {sorted(extensiones)}")
    
    # Ejemplo 3: Agregar nuevo lector dinámicamente
    print("\n3. EXTENSIBILIDAD (agregar nuevo formato):")
    print("-" * 50)
    
    class LectorEPUB(LectorArchivo):
        """Ejemplo de cómo agregar soporte para EPUB."""
        
        def puede_leer(self, extension: str) -> bool:
            return extension.lower() == '.epub'
        
        def leer(self, ruta: str) -> str:
            return f"[EPUB simulado] {ruta}"
    
    # Registrar el nuevo lector
    LectorFactory.registrar_lector(LectorEPUB())
    
    lector_epub = LectorFactory.crear_lector('libro.epub')
    print(f"  libro.epub → {type(lector_epub).__name__}")
    print("  ✅ Nuevo formato agregado sin modificar código existente!")
    
    # Ejemplo 4: Ventajas del patrón
    print("\n4. VENTAJAS DEL FACTORY PATTERN:")
    print("-" * 50)
    print("  ✅ Open/Closed: Abierto para extensión, cerrado para modificación")
    print("  ✅ Single Responsibility: Cada lector tiene una responsabilidad")
    print("  ✅ Extensible: Fácil agregar nuevos formatos")
    print("  ✅ Testeable: Cada lector se testea independientemente")
    print("  ✅ Mantenible: Bug en PDF solo afecta LectorPDF")
    
    print("\n" + "=" * 70)
    print("¡Implementación completa del Factory Pattern! 🏭✨")
    print("=" * 70)
