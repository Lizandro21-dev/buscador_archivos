"""
Utils - Módulo principal del Buscador de Archivos USB
Contiene la lógica de negocio, detección de unidades, búsqueda y gestión de archivos.
"""

import os
import sys
import string
import threading
import subprocess
import platform
from typing import List, Dict, Optional
import win32file

from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QPushButton, QLineEdit, QListWidget, QLabel, QMessageBox, 
    QFrame, QCompleter
)
from PyQt5.QtCore import Qt, QTimer, QStringListModel

from ventana import VentanaBase
from GestorHistorial import GestorHistorial


# ========== CONSTANTES ==========

# Extensiones de archivos que se pueden leer como texto
EXTENSIONES_TEXTO = {
    '.txt', '.log', '.csv', '.json', '.xml', '.html', 
    '.py', '.js', '.css', '.md', '.ini', '.conf'
}

# Extensiones de documentos que requieren librerías especiales
EXTENSIONES_DOCUMENTOS = {
    '.docx': 'docx',
    '.pptx': 'pptx', 
    '.xlsx': 'openpyxl',
    '.pdf': 'PyPDF2'
}

# Encodings a probar al leer archivos de texto
ENCODINGS_TEXTO = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']


# ========== FUNCIONES AUXILIARES ==========

def detectar_unidades_disponibles() -> List[Dict[str, str]]:
    """
    Detecta unidades extraíbles USB conectadas al sistema Windows.
    
    Proceso:
    1. Itera por todas las letras del alfabeto (A-Z)
    2. Verifica si existe la unidad y si es removible
    3. Excluye A:, B: (disqueteras) y C: (sistema)
    
    Returns:
        Lista de diccionarios con información de cada unidad USB:
        [{'texto': 'D:/', 'ruta': 'D:\\'}, ...]
    """
    unidades = []
    
    for letra in string.ascii_uppercase:
        unidad = f"{letra}:\\"
        
        # Saltar disqueteras y unidad del sistema
        if letra in ['A', 'B', 'C']:
            continue
        
        # Verificar si la unidad existe
        if not os.path.exists(unidad):
            continue
        
        try:
            # GetDriveType retorna:
            # 0 = Desconocido
            # 1 = No existe
            # 2 = Removible (USB)
            # 3 = Fijo (HDD)
            # 4 = Remoto (red)
            # 5 = CD-ROM
            # 6 = RAM Disk
            tipo_unidad = win32file.GetDriveType(unidad)
            
            if tipo_unidad == 2:  # Solo unidades removibles (USB)
                unidades.append({
                    'texto': f"{letra}:/", 
                    'ruta': unidad
                })
        except Exception as e:
            print(f"Error al detectar unidad {unidad}: {e}")
            continue
    
    return unidades


def leer_contenido_archivo(ruta: str) -> str:
    """
    Lee y extrae el contenido de texto de un archivo.
    
    Soporta múltiples formatos:
    - Archivos de texto plano (.txt, .log, .csv, etc.)
    - Documentos Word (.docx)
    - Presentaciones PowerPoint (.pptx)
    - Hojas de cálculo Excel (.xlsx)
    - Documentos PDF (.pdf)
    
    Args:
        ruta: Ruta completa del archivo a leer
        
    Returns:
        Contenido del archivo en minúsculas.
        String vacío si no se puede leer o el formato no es soportado.
    """
    try:
        extension = os.path.splitext(ruta)[1].lower()
        
        # Archivos de texto plano
        if extension in EXTENSIONES_TEXTO:
            return _leer_archivo_texto(ruta)
        
        # Documentos Word
        elif extension == '.docx':
            return _leer_docx(ruta)
        
        # Presentaciones PowerPoint
        elif extension == '.pptx':
            return _leer_pptx(ruta)
        
        # Hojas de cálculo Excel
        elif extension == '.xlsx':
            return _leer_xlsx(ruta)
        
        # Documentos PDF
        elif extension == '.pdf':
            return _leer_pdf(ruta)
        
        return ""
        
    except Exception as e:
        print(f"Error al leer archivo {ruta}: {e}")
        return ""


def _leer_archivo_texto(ruta: str) -> str:
    """
    Lee archivos de texto plano probando múltiples encodings.
    
    Args:
        ruta: Ruta del archivo de texto
        
    Returns:
        Contenido del archivo en minúsculas
    """
    for encoding in ENCODINGS_TEXTO:
        try:
            with open(ruta, 'r', encoding=encoding) as f:
                return f.read().lower()
        except (UnicodeDecodeError, IOError):
            continue
    return ""


def _leer_docx(ruta: str) -> str:
    """Lee contenido de archivos .docx usando python-docx."""
    try:
        import docx
        doc = docx.Document(ruta)
        texto = '\n'.join([p.text for p in doc.paragraphs])
        return texto.lower()
    except Exception:
        return ""


def _leer_pptx(ruta: str) -> str:
    """Lee contenido de archivos .pptx usando python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(ruta)
        texto = ''
        
        # Extraer texto de todas las diapositivas y formas
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + '\n'
        
        return texto.lower()
    except Exception:
        return ""


def _leer_xlsx(ruta: str) -> str:
    """Lee contenido de archivos .xlsx usando openpyxl."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(ruta, read_only=True, data_only=True)
        texto = ''
        
        # Extraer valores de todas las celdas de todas las hojas
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        texto += str(cell) + ' '
                texto += '\n'
        
        wb.close()
        return texto.lower()
    except Exception:
        return ""


def _leer_pdf(ruta: str) -> str:
    """Lee contenido de archivos .pdf usando PyPDF2."""
    try:
        import PyPDF2
        with open(ruta, 'rb') as f:
            pdf = PyPDF2.PdfReader(f)
            texto = ''
            for page in pdf.pages:
                texto += page.extract_text()
            return texto.lower()
    except Exception:
        return ""


# ========== CLASE PRINCIPAL ==========

class BuscadorArchivos(VentanaBase):
    """
    Clase principal del buscador de archivos USB.
    
    Funcionalidades:
    - Detección automática de unidades USB
    - Indexación de archivos
    - Búsqueda por nombre, extensión o contenido
    - Historial de búsquedas con autocompletado
    - Apertura de archivos con doble clic
    """
    
    # Constantes de configuración
    INTERVALO_DETECCION_USB = 2000  # ms (2 segundos)
    DELAY_BUSQUEDA_VIVO = 300       # ms (0.3 segundos)
    MAX_ARCHIVOS_MOSTRADOS = 500    # Límite para rendimiento
    
    # ELIMINADO: IDs de tipos de búsqueda (ya no se usan)
    # Ahora se busca por TODOS los criterios simultáneamente
    
    def __init__(self):
        """Inicializa el buscador de archivos."""
        super().__init__()
        
        # Estado de la aplicación
        self.unidad_seleccionada: Optional[str] = None
        self.resultados: List[Dict] = []
        self.todos_los_archivos: List[Dict] = []
        self.unidades_previas: List[Dict] = []
        # ELIMINADO: self.tipo_busqueda (ya no se usa)
        
        # Gestor de historial y autocompletado
        self.historial = GestorHistorial()
        self.completer: Optional[QCompleter] = None
        
        # Inicializar interfaz gráfica
        self.init_ui()
        self.cargar_unidades()
        self.configurar_autocompletado()
        
        # Timer para detectar cambios en USBs
        self.timer_usb = QTimer()
        self.timer_usb.timeout.connect(self.verificar_cambios_usb)
        self.timer_usb.start(self.INTERVALO_DETECCION_USB)
        
        # Timer para búsqueda en vivo (con delay)
        self.timer_autocompletar = QTimer()
        self.timer_autocompletar.setSingleShot(True)
        self.timer_autocompletar.timeout.connect(self.buscar_sugerencias)
    
    # ========== AUTOCOMPLETADO Y HISTORIAL ==========
    
    def configurar_autocompletado(self):
        """
        Configura el sistema de autocompletado con el historial de búsquedas.
        
        Features:
        - Búsqueda case-insensitive
        - Coincidencias parciales
        - Máximo 20 sugerencias visibles
        - Popup estilizado
        """
        self.completer = QCompleter()
        self.completer.setCaseSensitivity(Qt.CaseInsensitive)
        self.completer.setFilterMode(Qt.MatchContains)
        self.completer.setMaxVisibleItems(20)
        
        # Configurar estilo del popup de sugerencias
        popup = self.completer.popup()
        popup.setMinimumWidth(600)
        popup.setMinimumHeight(400)
        popup.setStyleSheet("""
            QListView {
                background-color: #2B313F;
                color: white;
                border: 3px solid #E32D64;
                border-radius: 12px;
                padding: 10px;
                font-size: 16px;
                font-family: 'Segoe UI', Arial;
            }
            QListView::item {
                padding: 12px;
                border-radius: 8px;
                margin: 3px;
                min-height: 30px;
            }
            QListView::item:selected {
                background-color: #E32D64;
                color: white;
                font-weight: bold;
            }
            QListView::item:hover {
                background-color: #4a4a4c;
            }
        """)
        
        self.search_input.setCompleter(self.completer)
        self.actualizar_completer()
        self.configurar_eventos_seleccion()
    
    def actualizar_completer(self):
        """Actualiza el modelo del completer con el historial actual."""
        model = QStringListModel()
        model.setStringList(self.historial.obtener_todos())
        self.completer.setModel(model)
    
    def configurar_eventos_seleccion(self):
        """
        Configura eventos personalizados para el campo de búsqueda.
        Permite mostrar el historial al hacer clic o enfocar.
        """
        # Guardar eventos originales
        self.search_input_click_original = self.search_input.mousePressEvent
        self.search_input_focus_original = self.search_input.focusInEvent
        
        # Reemplazar con eventos personalizados
        self.search_input.mousePressEvent = self.on_search_input_click
        self.search_input.focusInEvent = self.on_search_input_focus
    
    def on_search_input_click(self, event):
        """
        Maneja el evento de clic en el campo de búsqueda.
        Muestra el historial completo SIEMPRE, sin importar el texto.
        """
        # CAMBIO: SIEMPRE mostrar todo el historial al hacer clic
        try:
            model = QStringListModel()
            model.setStringList(self.historial.obtener_todos())
            self.completer.setModel(model)
            self.completer.complete()
        except Exception as e:
            print(f"Error al mostrar historial: {e}")
        
        # Ejecutar evento original para mantener funcionalidad normal
        self.search_input_click_original(event)
    
    def on_search_input_focus(self, event):
        """
        Maneja el evento cuando el campo obtiene el foco.
        Muestra el historial completo SIEMPRE.
        """
        self.search_input_focus_original(event)
        
        # CAMBIO: SIEMPRE mostrar todo el historial al enfocar
        try:
            model = QStringListModel()
            model.setStringList(self.historial.obtener_todos())
            self.completer.setModel(model)
        except Exception:
            pass
    
    def on_texto_cambiado(self):
        """
        Se ejecuta cuando cambia el texto en el campo de búsqueda.
        Actualiza sugerencias del historial e inicia búsqueda en vivo.
        """
        # CAMBIO: Siempre mostrar TODO el historial, no solo coincidencias
        # El QCompleter ya se encarga de filtrar internamente
        model = QStringListModel()
        model.setStringList(self.historial.obtener_todos())
        self.completer.setModel(model)
        
        # Iniciar búsqueda en vivo si hay unidad seleccionada
        self.iniciar_autocompletado()
    
    # ========== DETECCIÓN Y GESTIÓN DE USB ==========
    
    def cargar_unidades(self):
        """
        Detecta unidades USB disponibles y crea botones para cada una.
        Actualiza dinámicamente el panel izquierdo.
        """
        unidades = detectar_unidades_disponibles()
        
        # Limpiar botones anteriores
        self._limpiar_layout(self.units_layout)
        
        if not unidades:
            self._mostrar_mensaje_sin_unidades()
            return
        
        # Crear botón para cada unidad detectada
        for info in unidades:
            self._crear_boton_unidad(info)
    
    def _limpiar_layout(self, layout: QVBoxLayout):
        """Elimina todos los widgets de un layout."""
        for i in reversed(range(layout.count())):
            widget = layout.itemAt(i).widget()
            if widget:
                widget.setParent(None)
    
    def _mostrar_mensaje_sin_unidades(self):
        """Muestra un mensaje cuando no hay unidades USB detectadas."""
        lbl = QLabel("No hay unidades\nexternas detectadas")
        lbl.setAlignment(Qt.AlignCenter)
        lbl.setStyleSheet("color: white; font-size: 14px;")
        self.units_layout.addWidget(lbl)
    
    def _crear_boton_unidad(self, info: Dict[str, str]):
        """
        Crea un botón para una unidad USB.
        
        Args:
            info: Diccionario con 'texto' y 'ruta' de la unidad
        """
        btn = QPushButton(info['texto'])
        btn.setFixedHeight(60)
        btn.setStyleSheet("""
            QPushButton {
                background-color: #2B313F; 
                color: white;
                font-size: 20px;
                font-weight: bold;
                border: 2px solid #E32D64;
                border-radius: 15px;
            } 
            QPushButton:hover { 
                background-color: #4a4a4c; 
            }
            QPushButton:pressed {
                background-color: #E32D64;
            }
        """)
        btn.setCursor(Qt.PointingHandCursor)
        # Usar lambda para pasar el parámetro info correctamente
        btn.clicked.connect(lambda _, i=info: self.seleccionar_unidad(i))
        self.units_layout.addWidget(btn)
    
    def verificar_cambios_usb(self):
        """
        Verifica periódicamente si hay cambios en las USBs conectadas.
        Se ejecuta cada 2 segundos via timer.
        """
        actuales = detectar_unidades_disponibles()
        
        # Comparar solo las letras de unidad
        letras_actuales = {u['texto'] for u in actuales}
        letras_previas = {u['texto'] for u in self.unidades_previas}
        
        # Si hay cambios, recargar la interfaz
        if letras_actuales != letras_previas:
            self.cargar_unidades()
        
        self.unidades_previas = actuales
    
    # ========== SELECCIÓN E INDEXACIÓN DE UNIDADES ==========
    
    def seleccionar_unidad(self, unidad: Dict[str, str]):
        """
        Selecciona una unidad USB e inicia la indexación de sus archivos.
        
        Args:
            unidad: Diccionario con información de la unidad
        """
        self.unidad_seleccionada = unidad['ruta']
        
        # Mostrar mensaje de carga
        self.results_list.clear()
        self.results_list.addItem("")
        self.results_list.addItem(f"  Unidad seleccionada: {unidad['texto']}")
        self.results_list.addItem(f"  Ruta: {unidad['ruta']}")
        self.results_list.addItem("")
        self.results_list.addItem("  Indexando archivos, por favor espera...")
        
        # Indexar en segundo plano para no bloquear la interfaz
        threading.Thread(target=self.indexar_unidad, daemon=True).start()
    
    def indexar_unidad(self):
        """
        Recorre toda la unidad seleccionada e indexa todos sus archivos.
        Se ejecuta en un thread separado para no bloquear la UI.
        
        Guarda para cada archivo:
        - Nombre completo
        - Ruta completa
        - Extensión
        """
        self.todos_los_archivos = []
        
        if not self.unidad_seleccionada:
            return
        
        try:
            # Recorrer recursivamente todos los directorios
            for root, dirs, files in os.walk(self.unidad_seleccionada):
                for nombre_archivo in files:
                    ruta_completa = os.path.join(root, nombre_archivo)
                    
                    # Guardar información del archivo
                    self.todos_los_archivos.append({
                        'nombre': nombre_archivo,
                        'ruta': ruta_completa,
                        'extension': os.path.splitext(nombre_archivo)[1].lower()
                    })
            
            # Ordenar alfabéticamente por nombre
            self.todos_los_archivos.sort(key=lambda x: x['nombre'].lower())
            
            # Mostrar todos los archivos en la lista
            self.mostrar_todos_los_archivos()
            
        except Exception as e:
            self.alerta(f"Error al indexar unidad: {e}")
    
    def mostrar_todos_los_archivos(self):
        """
        Muestra todos los archivos indexados en la lista de resultados.
        Conecta el evento de doble clic para abrir archivos.
        """
        self.results_list.clear()
        self.resultados = self.todos_los_archivos
        
        # Conectar evento de doble clic (desconectar previo si existe)
        try:
            self.results_list.itemDoubleClicked.disconnect()
        except TypeError:
            pass
        self.results_list.itemDoubleClicked.connect(self.abrir_item)
        
        if not self.todos_los_archivos:
            self.results_list.addItem("")
            self.results_list.addItem("  No hay archivos en esta unidad")
            return
        
        # Mostrar encabezado con información
        self.results_list.addItem("")
        self.results_list.addItem(f"  Total: {len(self.todos_los_archivos)} archivos")
        self.results_list.addItem("")
        self.results_list.addItem("  Doble clic para abrir | Escribe para buscar")
        self.results_list.addItem("")
        self.results_list.addItem("  " + "=" * 80)
        
        # Mostrar archivos (limitar para mejor rendimiento)
        archivos_mostrados = self.todos_los_archivos[:self.MAX_ARCHIVOS_MOSTRADOS]
        for archivo in archivos_mostrados:
            self.results_list.addItem(f"  {archivo['nombre']}")
        
        # Indicar si hay más archivos
        archivos_restantes = len(self.todos_los_archivos) - self.MAX_ARCHIVOS_MOSTRADOS
        if archivos_restantes > 0:
            self.results_list.addItem("")
            self.results_list.addItem(f"  ... y {archivos_restantes} archivos más")
            self.results_list.addItem("  (Usa la búsqueda para filtrar)")
    
    # ========== TIPOS DE BÚSQUEDA ==========
    
    # MÉTODO ELIMINADO: Ya no se cambia el tipo de búsqueda
    # La búsqueda ahora es automática combinando todos los criterios
    """
    def cambiar_tipo_busqueda(self, boton):
        
        Cambia el tipo de búsqueda activo y actualiza la interfaz.
        
        Args:
            boton: Botón que fue clickeado
        
        self.tipo_busqueda = self.grupo_busqueda.id(boton)
        
        # Actualizar placeholder según el tipo de búsqueda
        placeholders = {
            self.BUSQUEDA_NOMBRE: "Buscar por nombre de archivo...",
            self.BUSQUEDA_EXTENSION: "Buscar por extensión (ej: .pdf, .txt)...",
            self.BUSQUEDA_CONTENIDO: "Buscar por contenido dentro de archivos..."
        }
        self.search_input.setPlaceholderText(placeholders.get(self.tipo_busqueda, ""))
        
        # Re-ejecutar búsqueda si hay texto ingresado
        if self.search_input.text().strip():
            self.buscar_sugerencias()
    """
    
    def iniciar_autocompletado(self):
        """
        Inicia el timer para búsqueda en vivo con delay.
        Evita hacer búsquedas por cada tecla presionada.
        """
        if not self.unidad_seleccionada or not self.todos_los_archivos:
            return
        
        # Reiniciar timer (espera 300ms desde última tecla)
        self.timer_autocompletar.start(self.DELAY_BUSQUEDA_VIVO)
    
    def BusquedaPor(self, texto: str) -> List[Dict]:
        """
        Busca archivos por TODOS los criterios: nombre, extensión y contenido.
        Combina los resultados de las tres búsquedas sin duplicados.
        
        Args:
            texto: Texto a buscar
            
        Returns:
            Lista de archivos únicos que coinciden con cualquiera de los criterios
        """
        # Usar un conjunto para evitar duplicados (basado en ruta de archivo)
        archivos_encontrados = {}
        
        # 1. Buscar por nombre
        resultados_nombre = self._buscar_por_nombre(texto)
        for archivo in resultados_nombre:
            archivos_encontrados[archivo['ruta']] = archivo
        
        # 2. Buscar por extensión
        resultados_extension = self._buscar_por_extension(texto)
        for archivo in resultados_extension:
            archivos_encontrados[archivo['ruta']] = archivo
        
        # 3. Buscar por contenido
        resultados_contenido = self._buscar_por_contenido(texto)
        for archivo in resultados_contenido:
            archivos_encontrados[archivo['ruta']] = archivo
        
        # Convertir el diccionario a lista
        return list(archivos_encontrados.values())
    
    def buscar_sugerencias(self):
        """
        Filtra archivos usando búsqueda combinada (nombre, extensión y contenido).
        Actualiza la lista de resultados en tiempo real.
        """
        texto = self.search_input.text().strip().lower()
        
        if not texto:
            # Sin texto, mostrar todos los archivos
            self.mostrar_todos_los_archivos()
            return
        
        # NUEVO: Usar BusquedaPor que combina todos los criterios
        coincidencias = self.BusquedaPor(texto)
        
        # Mostrar resultados filtrados
        self._mostrar_resultados(coincidencias, texto)
    
    def _buscar_por_nombre(self, texto: str) -> List[Dict]:
        """
        Busca archivos cuyo nombre contenga el texto (sin incluir extensión).
        
        Args:
            texto: Texto a buscar (en minúsculas)
            
        Returns:
            Lista de archivos que coinciden
        """
        coincidencias = []
        for archivo in self.todos_los_archivos:
            # Extraer nombre sin extensión
            nombre_sin_ext = os.path.splitext(archivo['nombre'])[0].lower()
            if texto in nombre_sin_ext:
                coincidencias.append(archivo)
        return coincidencias
    
    def _buscar_por_extension(self, texto: str) -> List[Dict]:
        """
        Busca archivos por extensión.
        
        Args:
            texto: Extensión a buscar (ej: ".pdf" o "pdf")
            
        Returns:
            Lista de archivos con esa extensión
        """
        # Normalizar: agregar punto si no lo tiene
        if not texto.startswith('.'):
            texto = '.' + texto
        
        coincidencias = []
        for archivo in self.todos_los_archivos:
            if archivo['extension'] == texto:
                coincidencias.append(archivo)
        return coincidencias
    
    def _buscar_por_contenido(self, texto: str) -> List[Dict]:
        """
        Busca archivos que contengan el texto dentro de su contenido.
        Esta búsqueda puede ser lenta con muchos archivos.
        
        Args:
            texto: Texto a buscar dentro de archivos
            
        Returns:
            Lista de archivos cuyo contenido contiene el texto
        """
        coincidencias = []
        total = len(self.todos_los_archivos)
        
        # Actualizar UI mostrando progreso
        self.results_list.clear()
        self.results_list.addItem("")
        self.results_list.addItem("  Buscando en contenido de archivos...")
        self.results_list.addItem("")
        
        for idx, archivo in enumerate(self.todos_los_archivos):
            # Mostrar progreso cada 10 archivos
            if idx % 10 == 0:
                self.results_list.item(2).setText(
                    f"  Progreso: {idx}/{total} archivos analizados..."
                )
                QApplication.processEvents()  # Actualizar UI
            
            # Leer contenido del archivo
            contenido = leer_contenido_archivo(archivo['ruta'])
            
            if texto in contenido:
                coincidencias.append(archivo)
        
        return coincidencias
    
    def _mostrar_resultados(self, coincidencias: List[Dict], texto_busqueda: str):
        """
        Muestra los resultados de búsqueda en la lista.
        
        Args:
            coincidencias: Lista de archivos que coinciden
            texto_busqueda: Texto que se buscó
        """
        self.results_list.clear()
        self.resultados = coincidencias
        
        # Conectar doble clic
        try:
            self.results_list.itemDoubleClicked.disconnect()
        except TypeError:
            pass
        self.results_list.itemDoubleClicked.connect(self.abrir_item)
        
        # Mostrar encabezado
        self.results_list.addItem("")
        self.results_list.addItem(f"  Búsqueda: '{texto_busqueda}'")
        self.results_list.addItem(f"  Resultados: {len(coincidencias)} archivo(s)")
        self.results_list.addItem("")
        self.results_list.addItem("  " + "=" * 80)
        
        if not coincidencias:
            self.results_list.addItem("")
            self.results_list.addItem("  No se encontraron archivos")
            return
        
        # Mostrar resultados
        for archivo in coincidencias[:self.MAX_ARCHIVOS_MOSTRADOS]:
            self.results_list.addItem(f"  {archivo['nombre']}")
        
        # Indicar si hay más resultados
        if len(coincidencias) > self.MAX_ARCHIVOS_MOSTRADOS:
            restantes = len(coincidencias) - self.MAX_ARCHIVOS_MOSTRADOS
            self.results_list.addItem("")
            self.results_list.addItem(f"  ... y {restantes} resultados más")
    
    # ========== BÚSQUEDA Y ACCIONES ==========
    
    def buscar_archivos(self):
        """
        Ejecuta la búsqueda cuando se presiona Enter o el botón Buscar.
        Agrega el término al historial.
        """
        texto = self.search_input.text().strip()
        
        if not texto:
            return
        
        if not self.unidad_seleccionada:
            self.alerta("Por favor, selecciona primero una unidad USB")
            return
        
        # Agregar al historial
        self.historial.agregar(texto)
        self.actualizar_completer()
        
        # Ejecutar búsqueda
        self.buscar_sugerencias()
    
    def abrir_item(self, item):
        """
        Abre el archivo seleccionado con doble clic.
        
        Args:
            item: Item de la lista que fue clickeado
        """
        texto_item = item.text().strip()
        
        # Buscar el archivo en los resultados actuales
        for archivo in self.resultados:
            if archivo['nombre'] == texto_item:
                self._abrir_archivo(archivo['ruta'])
                return
    
    def _abrir_archivo(self, ruta: str):
        """
        Abre un archivo con la aplicación predeterminada del sistema.
        
        Args:
            ruta: Ruta completa del archivo a abrir
        """
        try:
            if platform.system() == 'Windows':
                os.startfile(ruta)
            elif platform.system() == 'Darwin':  # macOS
                subprocess.call(['open', ruta])
            else:  # Linux
                subprocess.call(['xdg-open', ruta])
        except Exception as e:
            self.alerta(f"Error al abrir archivo: {e}")
    
    def limpiar_historial(self):
        """Limpia el historial de búsquedas con confirmación."""
        respuesta = QMessageBox.question(
            self,
            "Confirmar",
            "¿Deseas limpiar todo el historial de búsquedas?",
            QMessageBox.Yes | QMessageBox.No
        )
        
        if respuesta == QMessageBox.Yes:
            self.historial.limpiar()
            self.actualizar_completer()
            self.alerta("Historial limpiado exitosamente", QMessageBox.Information)
    
    def mostrar_mensaje_inicial(self):
        """Muestra el mensaje de bienvenida inicial."""
        self.results_list.clear()
        self.results_list.addItem("")
        self.results_list.addItem("  Bienvenido al Buscador de Archivos USB")
        self.results_list.addItem("")
        self.results_list.addItem("  Instrucciones:")
        self.results_list.addItem("  1. Selecciona una unidad USB del panel izquierdo")
        self.results_list.addItem("  2. Espera a que se indexen los archivos")
        self.results_list.addItem("  3. Escribe para buscar o explora todos los archivos")
        self.results_list.addItem("")
        self.results_list.addItem("  Búsqueda inteligente automática:")
        self.results_list.addItem("    • Busca simultáneamente en nombre, extensión y contenido")
        self.results_list.addItem("    • Obtén resultados más completos en una sola búsqueda")
        self.results_list.addItem("    • Ejemplo: 'python' encuentra archivos llamados 'python.txt',")
        self.results_list.addItem("      archivos .py y archivos que contengan la palabra 'python'")
        self.results_list.addItem("")
    
    def alerta(self, mensaje: str, tipo=QMessageBox.Warning):
        """
        Muestra un mensaje de alerta al usuario.
        
        Args:
            mensaje: Texto del mensaje
            tipo: Tipo de mensaje (Warning, Information, Critical)
        """
        msg = QMessageBox()
        msg.setIcon(tipo)
        msg.setText(mensaje)
        msg.setWindowTitle("Aviso")
        msg.exec_()


# ========== PUNTO DE ENTRADA ==========

def main():
    """Función principal para ejecutar la aplicación."""
    app = QApplication(sys.argv)
    ventana = BuscadorArchivos()
    ventana.show()
    sys.exit(app.exec_())


if __name__ == '__main__':
    main()